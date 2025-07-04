import os
import re
import json
from pptx import Presentation
from pptx.util import Inches
from apis.openai_api import OpenAIClient

default_config_path = "./config/config.json"

def get_config() -> dict:
    """Loads API configuration from JSON."""
    with open(default_config_path, 'r') as f:
        return json.load(f)

api_clients = {
    "openai": OpenAIClient
}

def clean_list(data):
    cleaned_data = []
    for item in data:
        title = item['title'].strip()
        content = item['content'].replace('\r', '')  # Remove carriage returns
        content = re.sub(r'\d+\.\s*', '', content)  # Remove numbered bullets
        content = content.strip() if content.strip() else ' '  # Add space if empty
        cleaned_data.append({'title': title, 'content': content})
    return cleaned_data

def get_generative_api_client(_generative_api_name, api_key, _generative_model_name):
    """Returns the AI client based on API name."""
    if _generative_api_name == "openai":
        return OpenAIClient(api_key, _generative_model_name)

def generate_ppt_content(topic, subject, class_name, num_slides):
    """Generates slide content using OpenAI before creating the PPT."""
    api_name = "openai"
    model_name = "gpt-4o-mini"
    config = get_config()

    final_prompt = f"""Create an outline for a slideshow presentation on the topic of {topic}. The presentation is for {subject}, Class {class_name}, and should be {num_slides} slides long.
        Ensure there are ONLY {num_slides} slides, including the title and thanks slide. The content should be informative, well-structured, and suitable for students.
        
        You are allowed to use the following slide types:
        Title Slide - (Title, Subtitle)
        Content Slide - (Title, Content)
        Thanks Slide - (Title)
        
        Put this tag before the Title Slide: [L_TS]
        Put this tag before the Content Slide: [L_CS]
        Put this tag before the Thanks Slide: [L_THS]
        
        Put this tag before the Title: [TITLE]
        Put this tag after the Title: [/TITLE]
        Put this tag before the Subtitle: [SUBTITLE]
        Put this tag after the Subtitle: [/SUBTITLE]
        Put this tag before the Content: [CONTENT]
        Put this tag after the Content: [/CONTENT]

        Put "[SLIDEBREAK]" after each slide 

        Example:
        [L_TS]
        [TITLE]Introduction to Algebra[/TITLE]
        [SUBTITLE]Mathematics - Class 10[/SUBTITLE]

        [SLIDEBREAK]

        [L_CS] 
        [TITLE]What is Algebra[/TITLE]
        [CONTENT]
        1. Algebra is a branch of mathematics dealing with symbols and the rules for manipulating them.
        2. It includes concepts like variables, equations, and expressions.
        3. It is used in various fields such as engineering, physics, and computer science.
        [/CONTENT]

        [SLIDEBREAK]

        The slides should be well-structured for easy comprehension. Do not include special characters (?, !, ., :, ) in the Title.
        Stick to the format exactly as described.
    """

    api_key = config.get('api_key')
    api_client = get_generative_api_client(api_name, api_key, model_name)
    presentation_content = api_client.generate(final_prompt)

    return parse_slide_content(presentation_content, topic, subject, class_name)

def parse_slide_content(response_text, topic, subject, class_name):
    """Parses AI-generated response into structured slide data and saves to a JSON file."""
    slides = []
    list_of_slides = response_text.split("[SLIDEBREAK]")

    def find_text_in_between_tags(text, start_tag, end_tag):
        """Extracts text between specified tags."""
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        if start_pos > -1 and end_pos > -1:
            return text[start_pos + len(start_tag):end_pos].strip()
        return ""

    for slide in list_of_slides:
        slide_type = next((tag for tag in ["[L_TS]", "[L_CS]", "[L_THS]"] if tag in slide), None)
        if slide_type:
            title = find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]")
            subtitle = find_text_in_between_tags(slide, "[SUBTITLE]", "[/SUBTITLE]")
            content = find_text_in_between_tags(slide, "[CONTENT]", "[/CONTENT]")

            slides.append({
                "type": slide_type,
                "title": title,
                "subtitle": subtitle if subtitle else "",
                "content": content if content else ""
            })

    # Create the folder structure: class_name/subject
    file_name = topic.replace(" ", "_")
    folder_path = os.path.join("Generated_Content", class_name, subject)
    os.makedirs(folder_path, exist_ok=True)

    # Save the slides to a JSON file within the folder
    file_path = os.path.join(folder_path, f'{file_name}.json')
    with open(file_path, 'w') as f:
        json.dump(slides, f, indent=4)
    return slides

def generate_final_ppt(topic, subject, class_name, slides):
    """Creates PPT using user-edited slides."""
    
    slides = clean_list(slides)
    print(slides)
    ppt = Presentation("theme0.pptx")

    def create_title_slide(title, subtitle):
        """Creates a title slide with subject and class details."""
        layout = ppt.slide_layouts[0]
        slide = ppt.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"Subject: {subject} | Class: {class_name}\n{subtitle}"

    def create_content_slide(title, content):
        """Creates a slide with title and content."""
        layout = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    def create_thanks_slide(title):
        """Creates a thank-you slide."""
        layout = ppt.slide_layouts[2]
        slide = ppt.slides.add_slide(layout)
        slide.shapes.title.text = title

    def delete_all_slides():
        """Deletes all slides from the PPT template before adding new ones."""
        for i in range(len(ppt.slides) - 1, -1, -1):
            r_id = ppt.slides._sldIdLst[i].rId
            ppt.part.drop_rel(r_id)
            del ppt.slides._sldIdLst[i]

    delete_all_slides()

    # ✅ Use slide order instead of relying on "type" (Fix for missing "type" error)
    for index, slide in enumerate(slides):
        title = slide.get("title", "Slide " + str(index + 1))  # Default to "Slide X" if missing
        content = slide.get("content", "")

        if index == 0:
            create_title_slide(title, content)  # First slide is the title slide
        elif index == len(slides) - 1:
            create_thanks_slide(title)  # Last slide is the thanks slide
        else:
            create_content_slide(title, content)  # Middle slides are content slides

    add_logo_to_slides(ppt, "logo.png")
    return ppt

def add_logo_to_slides(ppt, logo_path):
    """Adds a logo to all slides in the presentation."""
    logo_width = Inches(1.5)
    logo_height = logo_width * 98 / 183  # Maintain aspect ratio

    for slide in ppt.slides:
        left = ppt.slide_width - logo_width  # Adjust position
        top = ppt.slide_height - logo_height  # Adjust position
        slide.shapes.add_picture(logo_path, left, top, width=logo_width, height=logo_height)